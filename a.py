from bs4 import BeautifulSoup
import requests
import glob
import codecs
import os
import shutil
from pptx import Presentation
import frontmatter
import markdown
from jinja2 import Template, Environment, FileSystemLoader
import datetime
import os

def write_page(dir_path, filename, disp_text):
        os.makedirs(dir_path, exist_ok=True)
        with open(os.path.join(dir_path, filename), "w", encoding='utf-8') as f:
            f.write(disp_text)

def pptx(a,symbol):
    for b in a:
        print(b)
    prs = Presentation('./pptx/stock.pptx')
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text=='Company':
                    paragraph.text=a["company"]
                elif paragraph.text=='%Annual Dividend Yield':
                    paragraph.text=a["Forward_Annual_Dividend_Yield_4"]
                elif paragraph.text=='desc':
                    paragraph.text=a["disc"]
            for paragraph in shape.text_frame.paragraphs:
                print(paragraph.text)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for cell in shape.table.iter_cells():
                try:
                    cell.text=a[cell.text]
                except:
                    continue
            for cell in shape.table.iter_cells():
                print(cell.text)
    prs.save('./pptx/'+symbol+'.pptx')

def create_post(a):
    d_today = datetime.date.today()
    print(a["company"])
    dir_path = os.path.join('docs', 'blog')
    file_name = str(d_today) + '.md'
    template_path = "/template/post.md"
    env = Environment(loader=FileSystemLoader('.'), trim_blocks=False)
    template = env.get_template(template_path)
    disp_text = template.render(a=a,date=d_today)
    write_page(dir_path, file_name, disp_text)


def get_data(symbol):
    keyurl = 'https://finance.yahoo.com/quote/'+symbol
    r = requests.get(keyurl)
    soup = BeautifulSoup(r.text, features="lxml")
    table = soup.find(id="quote-summary")
    tr=table.find_all("td")
    td = [t.text.replace(" ","_") for t in tr]
    stockdict={td[i]:td[i+1] for i in range(0,len(td),2)}
    return stockdict

def static_data(symbol):
    keyurl = 'https://finance.yahoo.com/quote/'+symbol+"/key-statistics"
    r = requests.get(keyurl)
    soup = BeautifulSoup(r.text, features="lxml")
    tables=soup.find_all("td")
    td=[t.text.replace(" ","_") for t in tables]
    stockdict = {td[i]: td[i+1] for i in range(0, len(td), 2)}
    return stockdict

def div_his(symbol):
    keyurl = 'https://stocks.finance.yahoo.co.jp/us/profile/'+symbol
    r = requests.get(keyurl)
    soup = BeautifulSoup(r.text, features="lxml")
    tables = soup.find_all("tr")
    td = [t.text.replace(" ","_") for t in tables]
    disc={'disc':td[1].split("\n")[2],'name':td[0].split("\n")[1],'company':td[2].split("\n")[2],'adress':td[3].split("\n")[2],'from':td[5].split("\n")[2],'ceo':td[6].split("\n")[2],'sector':td[7].split("\n")[2],'market':td[8].split("\n")[2],"worker":td[9].split("\n")[2],"url":td[10].split("\n")[2]}
    return disc

def slide1(a):
    path=('C:\\Users\\tanimura\\Desktop\\stock\\ppt\\slides\\slide1.xml')
    print(a['Forward_Annual_Dividend_Yield_4'])
    text=[]
    with codecs.open(path,"r",'utf-8') as f:
            for row in f:
                text.append(row)
    text[1]=text[1].replace('%Annual Dividend Yield',a['Forward_Annual_Dividend_Yield_4'])
    with codecs.open(path,"w",'utf-8') as f:
        for txt in text:
            f.write(txt) 

if __name__ == '__main__':
    symbol = "A"
    a=get_data(symbol)
    b=div_his(symbol)
    c=static_data(symbol)
    a.update(b)
    a.update(c)
    pptx(a,symbol)
    create_post(a)
