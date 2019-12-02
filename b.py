from pptx import Presentation
prs = Presentation('C:\\Users\\tanimura\\Desktop\\stock.pptx')
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        for cell in shape.table.iter_cells():
            if cell.text == '%Payout Ratio%':
                cell.text='置換しました。'
        for cell in shape.table.iter_cells():
            print(cell.text)
print(text_runs)